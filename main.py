from flask import Flask, render_template, request, jsonify, session
import os
import io
from PyPDF2 import PdfReader
import fitz  # PyMuPDF for image extraction
from PIL import Image
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_google_genai import GoogleGenerativeAIEmbeddings
import google.generativeai as genai
from langchain.vectorstores import FAISS
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain.chains.question_answering import load_qa_chain
from langchain.prompts import PromptTemplate
from dotenv import load_dotenv
import docx  # for Word documents
import pandas as pd  # for Excel files
from pptx import Presentation  # for PowerPoint files
import datetime
import uuid
import threading
import json # For parsing LLM response
import traceback # For detailed error logging

app = Flask(__name__)
app.secret_key = os.urandom(24) # Needed for session management
app.config['UPLOAD_FOLDER'] = 'uploads'
app.config['MAX_CONTENT_LENGTH'] = 16 * 1024 * 1024  # 16MB max upload

# Create uploads folder if it doesn't exist
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)
os.makedirs('vector_stores', exist_ok=True) # Ensure vector store base exists

# Load environment variables
load_dotenv()
api_key = os.getenv("GOOGLE_API_KEY")
if not api_key:
    print("Error: GOOGLE_API_KEY not found in environment variables.")
    # Exit or handle gracefully if API key is mandatory
    exit("API Key is required.")
try:
    genai.configure(api_key=api_key)
except Exception as e:
    print(f"Error configuring GenerativeAI: {e}")
    exit("Failed to configure GenerativeAI.")


# Global dictionary to store processing status, messages, AND results
processing_status = {} # Structure: { session_id: { "status": ..., "progress": ..., "message": ..., "study_plan": [...] } }

# Lock for thread-safe access to processing_status (important if scaling later)
status_lock = threading.Lock()


def allowed_file(filename):
    """Check if file type is allowed"""
    allowed_extensions = {'pdf', 'docx', 'xlsx', 'xls', 'pptx', 'ppt', 'jpg', 'jpeg', 'png'}
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in allowed_extensions

# --- Text/Image Extraction Functions (keep as before) ---
def extract_text_from_word(doc_file):
    try:
        doc = docx.Document(doc_file)
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        return text
    except Exception as e:
        print(f"Error extracting text from Word: {e}")
        return ""

def extract_text_from_excel(excel_file):
    try:
        try:
            df = pd.read_excel(excel_file, engine='openpyxl')
        except Exception:
            df = pd.read_excel(excel_file)
        return df.to_string()
    except Exception as e:
        print(f"Error extracting text from Excel: {e}")
        return ""

def extract_text_from_ppt(ppt_file):
    try:
        prs = Presentation(ppt_file)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        print(f"Error extracting text from PowerPoint: {e}")
        return ""

def analyze_images_with_gemini(images):
    extracted_text = ""
    if not images:
        return ""
    try:
        model = genai.GenerativeModel("gemini-1.5-flash")
        prompt = "Analyze this image and describe its contents in detail, including any visible text, diagrams, or visual elements. Focus on extracting factual information presented."

        for i, img_data in enumerate(images):
            try:
                img = Image.open(io.BytesIO(img_data))
                if img.mode != 'RGB':
                   img = img.convert('RGB')
                response = model.generate_content([prompt, img])
                extracted_text += response.text + "\n\n" if response and response.text else ""
            except Exception as e:
                print(f"Error processing image {i+1} with Gemini: {str(e)}")
    except Exception as e:
        print(f"Error initializing Gemini model or during image analysis loop: {e}")

    return extracted_text

def get_text_chunks(text):
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=10000, chunk_overlap=1000)
    chunks = text_splitter.split_text(text)
    return chunks

def get_vector_store(text_chunks, session_id):
    try:
        embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
        vector_store = FAISS.from_texts(text_chunks, embedding=embeddings)

        session_dir = os.path.join('vector_stores', session_id)
        os.makedirs(session_dir, exist_ok=True)
        index_path = os.path.join(session_dir, "faiss_index")
        vector_store.save_local(index_path)
        print(f"Vector store saved to {index_path}")
    except Exception as e:
        print(f"Error creating or saving vector store: {e}")
        raise # Re-raise to be caught by the caller

# --- Study Plan Generation (keep as before) ---
def generate_study_plan(text, session_id):
    """Generates a study plan using Gemini based on the extracted text."""
    # IMPORTANT: This function now updates global status, use lock if needed
    # (Currently called only within the single processing thread, so lock might be overkill here,
    # but essential if this function were called from multiple places concurrently)
    try:
        # Update status message within the global dict
        with status_lock:
             if session_id in processing_status:
                 processing_status[session_id]["message"] = "Generating study plan structure..."

        model = genai.GenerativeModel("gemini-1.5-flash")

        max_text_length = 20000
        truncated_text = text[:max_text_length]

        prompt = f"""
        Analyze the following text extracted from one or more educational documents.
        Identify the main chapters, sections, or key topics suitable for a study plan checklist.
        Present the output as a JSON list of objects. Each object should represent a distinct study topic and have the following keys:
        - "id": A unique identifier string for the topic (e.g., "topic_1", "topic_2"). You MUST generate this ID.
        - "text": A concise string describing the topic title (e.g., "Chapter 1: Introduction", "1.2 Key Concepts", "Photosynthesis Process").

        Keep the topic titles relatively brief and actionable.
        Focus on the main structural elements (like headings, subheadings) or core concepts presented in the text.
        Aim for a reasonable number of topics (e.g., 10-30, depending on text length). Avoid overly granular topics unless they are distinct sub-sections.
        If the text is very short or unstructured, provide a list of the main themes or key takeaways.
        Ensure the output is ONLY the JSON list, enclosed in ```json ... ``` or just the raw list, without any preamble, explanation, or surrounding text.

        Text:
        ---
        {truncated_text}
        ---

        JSON Output:
        """

        response = model.generate_content(prompt)

        if response and response.text:
            json_text = response.text.strip().strip('```json').strip('```').strip()
            print(f"Raw Study Plan Response from LLM:\n{json_text[:500]}...")

            try:
                plan_items_raw = json.loads(json_text)
                study_plan = []
                if not isinstance(plan_items_raw, list):
                    print("Warning: LLM response for study plan was not a JSON list.")
                    raise json.JSONDecodeError("Response is not a list", json_text, 0)

                for i, item in enumerate(plan_items_raw):
                    if not isinstance(item, dict):
                         print(f"Warning: Skipping non-dictionary item in study plan list: {item}")
                         continue
                    topic_text = item.get("text", f"Unnamed Topic {i+1}").strip()
                    topic_id = item.get("id", "").strip()
                    if not topic_id:
                         topic_id = f"gen_topic_{uuid.uuid4()}"

                    if topic_text:
                        study_plan.append({
                            "id": topic_id,
                            "text": topic_text,
                            "completed": False
                        })
                print(f"Parsed Study Plan ({len(study_plan)} items): {study_plan[:5]}...")
                return study_plan

            except json.JSONDecodeError as json_e:
                print(f"Error decoding study plan JSON: {json_e}")
                print(f"Problematic JSON text: {json_text[:500]}...")
                # Fallback logic (keep as before)
                fallback_plan = []
                lines = json_text.split('\n')
                potential_topics = [line.strip() for line in lines if line.strip() and len(line.strip()) > 5]
                if potential_topics and any(l.startswith(('-', '*', '.')) or l[0].isdigit() for l in potential_topics):
                    print("JSON parsing failed, attempting fallback based on lines...")
                    for i, line in enumerate(potential_topics):
                         clean_line = line.lstrip('*-+. ').strip()
                         if clean_line:
                             fallback_plan.append({
                                 "id": f"fallback_topic_{uuid.uuid4()}",
                                 "text": clean_line,
                                 "completed": False
                             })
                    if fallback_plan:
                         print(f"Using fallback plan based on lines ({len(fallback_plan)} items).")
                         return fallback_plan
                return []
            except Exception as e:
                print(f"Error processing study plan items after JSON parsing: {e}")
                return []
        else:
            print("LLM did not return text for study plan.")
            return []

    except Exception as e:
        print(f"Error generating study plan: {e}")
        # Update status message on error
        with status_lock:
            if session_id in processing_status:
                 processing_status[session_id]["message"] = f"Study plan generation failed: {e}"
        return []


# --- MODIFIED process_documents --- Does NOT access session ---
def process_documents(files, session_id):
    """Process uploaded documents, create vector store, and generate study plan"""
    # This function runs in a background thread. DO NOT access `session` here.
    # Update status via the global `processing_status` dictionary using `status_lock`.
    start_time = datetime.datetime.now()
    print(f"[{start_time.strftime('%H:%M:%S')}] Background thread started for session {session_id}")

    # Local variables for processing
    current_progress = 0
    current_message = "Starting..."
    study_plan_result = [] # Store plan locally first
    success = True # Flag to track if process completed successfully

    def update_global_status(progress=None, message=None, status=None, plan=None, is_final=False):
        nonlocal current_progress, current_message # Allow modifying outer scope vars if needed
        with status_lock:
            if session_id not in processing_status:
                processing_status[session_id] = {} # Initialize if somehow missing

            if progress is not None:
                current_progress = progress
                processing_status[session_id]["progress"] = progress
            if message is not None:
                current_message = message
                processing_status[session_id]["message"] = message
            if status is not None:
                processing_status[session_id]["status"] = status
            if is_final and plan is not None: # Only store final plan
                 processing_status[session_id]["study_plan"] = plan

    try:
        update_global_status(progress=0, message="Initializing document processing...", status="processing")

        total_files = len(files)
        text = ""
        images = []

        # --- Stage 1: File Reading (40%) ---
        for i, file_path in enumerate(files):
            progress = int((i / total_files) * 40)
            filename = os.path.basename(file_path)
            update_global_status(progress=progress, message=f"Processing file {i+1}/{total_files}: {filename}")
            print(f"[{datetime.datetime.now().strftime('%H:%M:%S')}] Processing {filename}")
            # ... (rest of file reading logic - identical to before) ...
            file_extension = os.path.splitext(file_path)[1].lower()
            try:
                if file_extension == '.pdf':
                    pdf_reader = PdfReader(file_path)
                    for page in pdf_reader.pages: text += (page.extract_text() or "") + "\n"
                    doc_pdf = fitz.open(file_path)
                    for page in doc_pdf:
                        for img in page.get_images(full=True):
                            base_image = doc_pdf.extract_image(img[0])
                            images.append(base_image["image"])
                    doc_pdf.close()
                elif file_extension == '.docx': text += extract_text_from_word(file_path) + "\n"
                elif file_extension in ['.xlsx', '.xls']: text += extract_text_from_excel(file_path) + "\n"
                elif file_extension in ['.pptx', '.ppt']: text += extract_text_from_ppt(file_path) + "\n"
                elif file_extension in ['.jpg', '.jpeg', '.png']:
                     with open(file_path, 'rb') as img_file: images.append(img_file.read())
            except Exception as e:
                 print(f"Error processing {filename} in thread: {e}")
                 update_global_status(message=f"Error on {filename}: {e}. Skipping.") # Update message but continue

        # --- Stage 2: Image Analysis (15%) ---
        update_global_status(progress=40, message=f"Analyzing {len(images)} image(s)..." if images else "No images to analyze")
        image_text = ""
        if images:
            print(f"[{datetime.datetime.now().strftime('%H:%M:%S')}] Analyzing {len(images)} images...")
            image_text = analyze_images_with_gemini(images)
        update_global_status(progress=55)
        combined_text = text + "\n" + image_text

        # --- Stage 3: Study Plan Generation (15%) ---
        update_global_status(progress=55, message="Generating study plan...")
        print(f"[{datetime.datetime.now().strftime('%H:%M:%S')}] Generating study plan...")
        if combined_text.strip():
             # Call generate_study_plan (which might update status internally too)
             study_plan_result = generate_study_plan(combined_text, session_id)
             print(f"[{datetime.datetime.now().strftime('%H:%M:%S')}] Study plan generated with {len(study_plan_result)} items.")
        else:
             print(f"[{datetime.datetime.now().strftime('%H:%M:%S')}] Skipping study plan generation (no text).")
             study_plan_result = []
        update_global_status(progress=70) # Mark plan generation step complete

        # --- Stage 4: Text Chunking (15%) ---
        update_global_status(progress=70, message="Chunking text for Q&A...")
        print(f"[{datetime.datetime.now().strftime('%H:%M:%S')}] Chunking text...")
        text_chunks = []
        if combined_text.strip():
             text_chunks = get_text_chunks(combined_text)
        print(f"[{datetime.datetime.now().strftime('%H:%M:%S')}] Text chunking complete ({len(text_chunks)} chunks).")
        update_global_status(progress=85)

        # --- Stage 5: Vector Store Creation (15%) ---
        update_global_status(progress=85, message="Creating vector store...")
        print(f"[{datetime.datetime.now().strftime('%H:%M:%S')}] Creating vector store...")
        if text_chunks:
             get_vector_store(text_chunks, session_id)
        else:
             print(f"[{datetime.datetime.now().strftime('%H:%M:%S')}] Skipping vector store creation (no text chunks).")
             session_dir = os.path.join('vector_stores', session_id)
             os.makedirs(session_dir, exist_ok=True) # Still create dir
        print(f"[{datetime.datetime.now().strftime('%H:%M:%S')}] Vector store step complete.")

        # --- Final Success Update ---
        update_global_status(progress=100,
                             message="Processing complete! Ready.",
                             status="completed",
                             plan=study_plan_result, # Store the final plan now
                             is_final=True)
        success = True

    except Exception as e:
        success = False
        print(f"!!! Critical Error in process_documents thread for session {session_id}: {e}")
        traceback.print_exc() # Print full traceback for debugging
        error_message = f"Processing failed unexpectedly: {str(e)}"
        # --- Final Failure Update ---
        update_global_status(progress=current_progress, # Keep last known progress
                             message=error_message,
                             status="failed",
                             plan=[], # Ensure plan is empty on failure
                             is_final=True)

    finally:
        end_time = datetime.datetime.now()
        status_str = "Finished" if success else "Failed"
        print(f"[{end_time.strftime('%H:%M:%S')}] Background thread {status_str} for session {session_id}. Duration: {end_time - start_time}")


# --- Flask Routes ---

@app.route('/')
def index():
    # Ensure session_id exists
    if 'session_id' not in session:
        session['session_id'] = str(uuid.uuid4())
        print(f"New session created: {session['session_id']}")
    else:
        print(f"Existing session found: {session['session_id']}")

    # Initialize status for new/existing session if not already processing
    session_id = session['session_id']
    with status_lock:
        if session_id not in processing_status or processing_status[session_id].get("status") not in ["processing", "starting"]:
             processing_status[session_id] = {
                 "status": "idle",
                 "progress": 0,
                 "message": "Ready to process documents.",
                 "study_plan": session.get('study_plan', []) # Try to load from session initially? Or clear? Let's clear.
                 #"study_plan": [] # Clear plan on page load
             }
             # Clear potentially stale session plan if resetting status
             # session.pop('study_plan', None) # Let's avoid modifying session here, handle plan via processing_status

    return render_template('index.html')


@app.route('/upload', methods=['POST'])
def upload_files():
    if 'files[]' not in request.files:
        return jsonify({'error': 'No files part in the request'}), 400

    files = request.files.getlist('files[]')
    if not files or all(file.filename == '' for file in files):
        return jsonify({'error': 'No files selected'}), 400

    # Ensure session_id exists (should be set by '/' route)
    if 'session_id' not in session:
        # This shouldn't normally happen if '/' is visited first
        session['session_id'] = str(uuid.uuid4())
        print(f"Warning: Session ID created during upload for {session['session_id']}")
    session_id = session['session_id']

    user_upload_dir = os.path.join(app.config['UPLOAD_FOLDER'], session_id)
    os.makedirs(user_upload_dir, exist_ok=True)

    saved_files = []
    invalid_files = []
    # --- (File saving logic - identical to before) ---
    for file in files:
        if file and allowed_file(file.filename):
            filename = file.filename
            file_path = os.path.join(user_upload_dir, filename)
            try:
                file.save(file_path)
                saved_files.append(file_path)
            except Exception as e:
                print(f"Error saving file {filename}: {e}")
                return jsonify({'error': f'Error saving file: {str(e)}'}), 500
        elif file and file.filename != '':
            invalid_files.append(file.filename)

    if not saved_files:
         # Handle cases with only invalid files or no files saved
         error_msg = 'No valid files were successfully saved.'
         if invalid_files:
             error_msg += f' Invalid files skipped: {", ".join(invalid_files)}'
         return jsonify({'error': error_msg}), 400


    # --- Initialize processing status BEFORE starting thread ---
    with status_lock:
        processing_status[session_id] = {
            "status": "starting",
            "progress": 0,
            "message": "Initializing processing...",
            "study_plan": [] # Initialize empty plan
        }

    # --- Start background thread ---
    processing_thread = threading.Thread(target=process_documents, args=(saved_files, session_id), daemon=True)
    # Daemon=True allows app to exit even if threads are running (use carefully)
    processing_thread.start()
    print(f"Started background processing thread for session {session_id}")

    file_names = [os.path.basename(f) for f in saved_files]
    return jsonify({
        'message': f'Successfully uploaded {len(saved_files)} files. Processing started.',
        'files': file_names,
        'session_id': session_id,
        'warnings': f'Skipped invalid files: {", ".join(invalid_files)}' if invalid_files else None
    })


@app.route('/status', methods=['GET'])
def get_status():
    session_id = session.get('session_id')
    if not session_id:
        return jsonify({
            'status': 'no_session', 'progress': 0,
            'message': 'Session not found. Please reload.'
        }), 404

    with status_lock:
        # Get a copy to avoid potential modification during jsonify
        status_data = processing_status.get(session_id, {}).copy()

    # Don't send the full study plan in the status response for efficiency
    status_data.pop('study_plan', None)

    if not status_data: # If session_id existed but no status entry (shouldn't happen often)
        status_data = {'status': 'idle', 'progress': 0, 'message': 'Ready to process.'}

    # Ensure essential keys exist
    status_data.setdefault('status', 'unknown')
    status_data.setdefault('progress', 0)
    status_data.setdefault('message', '')

    # print(f"Reporting status for {session_id}: {status_data}") # Verbose
    return jsonify(status_data)


# --- MODIFIED /study_plan --- Retrieves from global status ---
@app.route('/study_plan', methods=['GET'])
def get_study_plan():
    """Endpoint to fetch the generated study plan for the current session."""
    session_id = session.get('session_id')
    if not session_id:
        return jsonify({'error': 'Session not found'}), 404

    with status_lock:
        current_status = processing_status.get(session_id, {})
        status = current_status.get('status', 'unknown')

        # Only return plan if processing is completed
        if status == 'completed':
            study_plan = current_status.get('study_plan', [])
            print(f"Returning study plan for completed session {session_id}: {len(study_plan)} items")
            return jsonify({'study_plan': study_plan})
        elif status == 'processing' or status == 'starting':
             print(f"Study plan requested but processing not complete for session {session_id} (Status: {status})")
             return jsonify({'study_plan': [], 'message': 'Processing not complete.'}) # Return empty list while processing
        else: # idle, failed, unknown, etc.
             print(f"Study plan requested but no completed plan found for session {session_id} (Status: {status})")
             return jsonify({'study_plan': []}) # Return empty list if not completed or failed


# --- MODIFIED /update_topic_status --- Updates global status ---
@app.route('/update_topic_status', methods=['POST'])
def update_topic_status():
    """Endpoint to update the completion status of a study plan topic."""
    session_id = session.get('session_id')
    if not session_id:
        return jsonify({'error': 'Session not found'}), 404

    data = request.json
    topic_id = data.get('topic_id')
    completed = data.get('completed') # Should be boolean

    if topic_id is None or completed is None:
        return jsonify({'error': 'Missing topic_id or completed status'}), 400

    updated = False
    with status_lock:
        if session_id in processing_status and processing_status[session_id].get('status') == 'completed':
            study_plan = processing_status[session_id].get('study_plan', [])
            for topic in study_plan:
                if topic.get('id') == topic_id:
                    topic['completed'] = bool(completed)
                    updated = True
                    # No need to save back explicitly, modifying the dict in place
                    break
            if updated:
                 print(f"Updated topic '{topic_id}' status to {completed} in global status for session {session_id}")
                 # Optionally: Persist this change somewhere more permanent if needed beyond server restart
            else:
                 print(f"Topic '{topic_id}' not found in completed plan for session {session_id}")
        else:
             print(f"Cannot update topic status: Processing not complete or session '{session_id}' not found.")
             return jsonify({'error': 'Processing not complete or plan not available'}), 400 # Or maybe 409 Conflict

    if updated:
        return jsonify({'success': True, 'message': 'Topic status updated'})
    else:
        # This case means topic_id wasn't found in the plan
        return jsonify({'error': 'Topic not found in study plan'}), 404

# --- /query, /export, get_answer (Keep as before, they don't use session in problematic ways) ---
@app.route('/query', methods=['POST'])
def query_documents():
    # ... (keep existing implementation) ...
    # It correctly checks processing_status before calling get_answer
    data = request.json
    question = data.get('question')
    answer_type = data.get('answer_type', 'Explanatory')
    session_id = session.get('session_id')

    print(f"Received query: question='{question[:50]}...', answer_type={answer_type}, session_id={session_id}")

    if not question: return jsonify({'error': 'No question provided'}), 400
    if not session_id: return jsonify({'error': 'No session ID found. Please reload.'}), 400

    with status_lock: # Check status safely
        current_status_dict = processing_status.get(session_id, {})
        current_status = current_status_dict.get('status')

    if current_status != 'completed':
        print(f"Error: Session {session_id} processing status is '{current_status or 'unknown'}'")
        return jsonify({'error': 'Document processing is not complete or has failed.'}), 400

    try:
        print("Attempting to get answer...")
        answer = get_answer(question, answer_type, session_id)
        print(f"Answer generated successfully.")
        return jsonify({'answer': answer})
    except FileNotFoundError as fnf_error:
         print(f"Error during get_answer (FileNotFound): {str(fnf_error)}")
         return jsonify({'error': 'Could not load document data. Please try processing again.'}), 500
    except Exception as e:
        print(f"Error during get_answer: {str(e)}", exc_info=True)
        return jsonify({'error': f'Error processing query: {str(e)}'}), 500

@app.route('/export', methods=['POST'])
def export_chat():
    # ... (keep existing implementation) ...
    data = request.json
    chat_messages = data.get('messages', [])
    format_type = data.get('format', 'markdown')

    if not chat_messages:
        return jsonify({'error': 'No chat messages to export'}), 400

    timestamp = datetime.datetime.now().strftime('%Y%m%d_%H%M%S')
    filename = f"docchat_export_{timestamp}"
    content = "" # Generate content based on format_type (HTML, MD, TXT)
    # ... (keep content generation logic) ...
    if format_type == 'html':
        content = """<!DOCTYPE html><html lang="en"><head><meta charset="UTF-8"><title>Document Chat Export</title><style>body { font-family: sans-serif; line-height: 1.5; max-width: 800px; margin: 20px auto; padding: 15px; border: 1px solid #ddd; }.message { margin-bottom: 1em; padding: 0.8em; border-radius: 5px; }.user { background-color: #e1f5fe; border-left: 4px solid #4fc3f7; }.ai { background-color: #e8f5e9; border-left: 4px solid #81c784; }.system { background-color: #fff3e0; border-left: 4px solid #ffb74d; font-style: italic; }strong { font-weight: bold; }p { margin: 0 0 0.5em 0; }p:last-child { margin-bottom: 0; }</style></head><body><h1>Document Chat Export</h1><p>Exported on: """ + datetime.datetime.now().strftime('%Y-%m-%d %H:%M:%S') + """</p><hr>"""
        for msg in chat_messages:
            role = msg.get('role', 'unknown')
            message_text = msg.get('content', '').replace('<', '&lt;').replace('>', '&gt;')
            message_html = message_text.replace('\n', '<br>')
            if role == 'user': content += f"<div class='message user'><strong>You:</strong><br>{message_html}</div>\n"
            elif role == 'ai': content += f"<div class='message ai'><strong>AI:</strong><br>{message_html}</div>\n"
            elif role == 'system': content += f"<div class='message system'><strong>System:</strong><br>{message_html}</div>\n"
        content += "</body></html>"
        filename += ".html"
    elif format_type == 'markdown':
        content = f"# Document Chat Export\n\n**Exported on:** {datetime.datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n\n---\n\n"
        for msg in chat_messages:
            role = msg.get('role', 'unknown')
            message = msg.get('content', '')
            if role == 'user': content += f"**You:**\n{message}\n\n---\n\n"
            elif role == 'ai': content += f"**AI:**\n{message}\n\n---\n\n"
            elif role == 'system': content += f"***System:***\n*{message}*\n\n---\n\n"
        filename += ".md"
    else:  # Plain text
        content = f"Document Chat Export\nExported on: {datetime.datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n\n====================\n\n"
        for msg in chat_messages:
            role = msg.get('role', 'unknown')
            message = msg.get('content', '')
            if role == 'user': content += f"You:\n{message}\n\n"
            elif role == 'ai': content += f"AI:\n{message}\n\n"
            elif role == 'system': content += f"System:\n{message}\n\n"
            content += "====================\n\n"
        filename += ".txt"

    return jsonify({'content': content, 'filename': filename})

def get_answer(question, answer_type, session_id):
    # ... (keep existing implementation) ...
    # This function is called within the /query request context, so it's fine.
    try:
        print(f"Getting answer for question: '{question[:50]}...', type: {answer_type}, session: {session_id}")
        prompt_templates = { # Prompts omitted for brevity - keep your existing ones
            "Explanatory": "...", "Concise": "...", "Cheat Sheet": "...", "Exam-Ready": "..."
        }
        prompt_templates = {
            "Explanatory": """
            You are a helpful study assistant. Answer the question based *only* on the provided context, which includes text extracted from documents and descriptions of images found within them.
            Provide detailed explanations, elaborate on concepts, and use examples if found in the context.
            Structure your answer clearly. If the context is insufficient or doesn't contain the answer, state that clearly, for example: "Based on the provided documents, the answer is not available." Do not make up information.

            Context:
            {context}

            Question:
            {question}

            Detailed Answer:
            """,
            "Concise": """
            You are a succinct study assistant. Answer the question briefly and directly based *only* on the provided context (text and image descriptions).
            Focus on the key points relevant to the question. Use short sentences.
            If the answer is not in the context, state: "The answer is not available in the provided documents." Do not guess or add external information.

            Context:
            {context}

            Question:
            {question}

            Concise Answer:
            """,
            "Cheat Sheet": """
            You are an organized study assistant. Extract the key information relevant to the question from the provided context (text and image descriptions) and present it in a cheat sheet format.
            Use bullet points, numbered lists, or key definitions. Focus on facts, summaries, or steps found in the context.
            If the context does not contain relevant information for a cheat sheet on this question, state: "No specific details for a cheat sheet on this topic were found in the documents." Do not invent information.

            Context:
            {context}

            Question:
            {question}

            Cheat Sheet Format Answer:
            """,
            "Exam-Ready": """
            You are a knowledgeable study assistant preparing a student for an exam. Answer the question thoroughly and in a well-structured format, based *only* on the provided context (text and image descriptions).
            Imagine you are writing an answer that would receive full marks. Use a clear introduction, body (with supporting points from the context), and conclusion if appropriate. Use precise language found in the context where possible.
            If the information needed to answer the question comprehensively is not in the context, state that clearly, for example: "The provided documents do not contain sufficient information to fully answer this question in an exam-ready format." Do not add external knowledge.

            Context:
            {context}

            Question:
            {question}

            Exam-Ready Answer:
            """
        }

        session_dir = os.path.join('vector_stores', session_id)
        vector_store_path = os.path.join(session_dir, "faiss_index")
        print(f"Attempting to load vector store from: {vector_store_path}")

        # Check for the actual index file, not just the directory
        index_file_path = os.path.join(vector_store_path, "index.faiss")
        if not os.path.exists(index_file_path):
             print(f"Vector store index file not found at: {index_file_path}")
             return "No document content was processed or the vector store is missing. Please process documents first."

        try:
            embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
            new_db = FAISS.load_local(vector_store_path, embeddings, allow_dangerous_deserialization=True)
            print("Vector store loaded successfully.")
        except Exception as load_error:
            print(f"Error loading FAISS index from {vector_store_path}: {load_error}")
            raise FileNotFoundError(f"Failed to load vector store from {vector_store_path}. Details: {load_error}")

        print(f"Performing similarity search for: '{question[:50]}...'")
        docs = new_db.similarity_search(question, k=5)
        if not docs:
            print("No relevant document chunks found.")
            return "No relevant information found in the processed documents for your question."

        print(f"Found {len(docs)} relevant document chunks.")

        selected_template = prompt_templates.get(answer_type, prompt_templates["Explanatory"])
        prompt = PromptTemplate(template=selected_template, input_variables=["context", "question"])
        model = ChatGoogleGenerativeAI(model="gemini-1.5-flash-latest", temperature=0.4, convert_system_message_to_human=True)
        chain = load_qa_chain(model, chain_type="stuff", prompt=prompt, verbose=False)

        print("Running QA chain...")
        response = chain({"input_documents": docs, "question": question}, return_only_outputs=True)
        final_answer = response.get("output_text", "Error: Could not generate an answer from the model.")
        return final_answer

    except FileNotFoundError as e:
         print(f"Vector store not found error in get_answer: {e}")
         raise e
    except Exception as e:
        print(f"Unexpected error in get_answer: {str(e)}", exc_info=True)
        return f"An unexpected error occurred while generating the answer: {str(e)}"


if __name__ == '__main__':
    print("Starting Flask app...")
    # Use host='0.0.0.0' to make accessible on network, default port 5000
    # Set debug=False for production environments
    app.run(debug=True, host='0.0.0.0', port=5000)